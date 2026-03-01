import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from services.scanner import FileScanner


class TestFileScanner:

    def test_text_extension_recognized(self, tmp_path):
        for name in ['main.py', 'app.ts', 'readme.md']:
            (tmp_path / name).write_text('content')

        results = FileScanner().scan_directory(str(tmp_path))

        assert len(results) == 3

    def test_extract_text_file(self, tmp_path):
        f = tmp_path / 'hello.txt'
        f.write_text('hello world')

        assert FileScanner().extract_text(str(f)) == 'hello world'

    def test_unsupported_extension_returns_none(self, tmp_path):
        f = tmp_path / 'file.xyz'
        f.write_text('data')

        assert FileScanner().extract_text(str(f)) is None

    def test_image_files_skipped_in_scan(self, tmp_path):
        for name in ['photo.png', 'pic.jpg', 'anim.gif']:
            (tmp_path / name).write_text('')

        assert FileScanner().scan_directory(str(tmp_path)) == []

    def test_skip_dirs_filtered(self, tmp_path):
        (tmp_path / 'node_modules').mkdir()
        (tmp_path / 'node_modules' / 'foo.js').write_text('x')
        (tmp_path / 'src').mkdir()
        (tmp_path / 'src' / 'bar.js').write_text('x')

        results = FileScanner().scan_directory(str(tmp_path))

        assert len(results) == 1
        assert results[0].endswith('bar.js')

    def test_skip_files_filtered(self, tmp_path):
        (tmp_path / '.DS_Store').write_text('')
        (tmp_path / 'readme.md').write_text('hi')

        results = FileScanner().scan_directory(str(tmp_path))

        assert len(results) == 1
        assert results[0].endswith('readme.md')

    def test_gitignore_respected(self, tmp_path):
        (tmp_path / '.gitignore').write_text('*.log\n')
        (tmp_path / 'debug.log').write_text('log data')
        (tmp_path / 'app.py').write_text('code')

        results = FileScanner().scan_directory(str(tmp_path))

        assert len(results) == 1
        assert results[0].endswith('app.py')

    def test_nested_gitignore(self, tmp_path):
        sub = tmp_path / 'sub'
        sub.mkdir()
        (sub / '.gitignore').write_text('secret.txt\n')
        (sub / 'secret.txt').write_text('secret')
        (sub / 'code.py').write_text('code')

        results = FileScanner().scan_directory(str(tmp_path))

        assert len(results) == 1
        assert results[0].endswith('code.py')

    def test_estimate_folder_size_only_supported(self, tmp_path):
        (tmp_path / 'app.py').write_text('code')
        (tmp_path / 'photo.png').write_text('img')
        (tmp_path / 'file.xyz').write_text('unknown')

        total_bytes, count, files = FileScanner().estimate_folder_size(str(tmp_path))

        assert count == 1
        assert files[0].endswith('app.py')
        assert total_bytes == len('code')

    def test_docx_extraction(self, tmp_path):
        from docx import Document
        path = str(tmp_path / 'test.docx')
        doc = Document()
        doc.add_paragraph('Hello from docx')
        doc.save(path)

        assert 'Hello from docx' in FileScanner().extract_text(path)

    def test_xlsx_extraction(self, tmp_path):
        from openpyxl import Workbook
        path = str(tmp_path / 'test.xlsx')
        wb = Workbook()
        ws = wb.active
        ws.append(['Name', 'Age'])
        ws.append(['Alice', 30])
        wb.save(path)

        text = FileScanner().extract_text(path)
        assert 'Name' in text
        assert 'Alice' in text

    def test_pptx_extraction(self, tmp_path):
        from pptx import Presentation
        path = str(tmp_path / 'test.pptx')
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = 'Slide Title'
        prs.save(path)

        assert 'Slide Title' in FileScanner().extract_text(path)

    def test_extract_text_with_meta_pdf(self, tmp_path):
        from fpdf import FPDF

        path = str(tmp_path / 'test.pdf')
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Helvetica', size=12)
        pdf.cell(text='Page one content')
        pdf.add_page()
        pdf.cell(text='Page two content')
        pdf.output(path)

        result = FileScanner().extract_text_with_meta(path)

        assert len(result) == 2
        assert 'Page one' in result[0][0]
        assert result[0][1] == {'page': 1}
        assert 'Page two' in result[1][0]
        assert result[1][1] == {'page': 2}

    def test_extract_text_with_meta_pptx(self, tmp_path):
        from pptx import Presentation
        path = str(tmp_path / 'test.pptx')
        prs = Presentation()
        for i in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f'Slide {i + 1}'
        prs.save(path)

        result = FileScanner().extract_text_with_meta(path)

        assert len(result) == 2
        assert 'Slide 1' in result[0][0]
        assert result[0][1] == {'slide': 1}
        assert 'Slide 2' in result[1][0]
        assert result[1][1] == {'slide': 2}

    def test_extract_text_with_meta_xlsx(self, tmp_path):
        from openpyxl import Workbook
        path = str(tmp_path / 'test.xlsx')
        wb = Workbook()
        ws1 = wb.active
        ws1.title = 'Sales'
        ws1.append(['Product', 'Revenue'])
        ws2 = wb.create_sheet('Costs')
        ws2.append(['Item', 'Amount'])
        wb.save(path)

        result = FileScanner().extract_text_with_meta(path)

        assert len(result) == 2
        assert result[0][1] == {'sheet': 'Sales'}
        assert 'Product' in result[0][0]
        assert result[1][1] == {'sheet': 'Costs'}
        assert 'Item' in result[1][0]

    def test_extract_text_with_meta_text(self, tmp_path):
        f = tmp_path / 'hello.txt'
        f.write_text('hello world')

        result = FileScanner().extract_text_with_meta(str(f))

        assert len(result) == 1
        assert result[0][0] == 'hello world'
        assert result[0][1] == {'type': 'text'}
