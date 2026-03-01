import os
import pathspec
from PyPDF2 import PdfReader


class FileScanner:
    TEXT_EXTENSIONS = {
        '.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv', '.log',
        '.ts', '.tsx', '.jsx', '.go', '.rs', '.java', '.rb', '.php', '.swift', '.kt',
        '.c', '.cpp', '.h', '.sh', '.yaml', '.yml', '.toml', '.sql', '.r', '.scala',
        '.dart', '.lua', '.pl', '.ex', '.exs', '.hs', '.ml', '.vim', '.dockerfile',
        '.makefile', '.gradle', '.bat', '.ps1', '.env',
        '.rtf', '.tex', '.rst', '.org', '.adoc',
    }

    OFFICE_EXTENSIONS = {'.docx', '.xlsx', '.pptx'}

    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.svg', '.webp', '.ico', '.bmp'}

    SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | OFFICE_EXTENSIONS | {'.pdf'}

    SKIP_DIRS = {
        'node_modules', '__pycache__', '.git', '.venv', 'venv', '.env', 'dist',
        'build', '.next', '.nuxt', '.cache', '.idea', '.vscode', 'vendor',
        'target', '.gradle',
    }

    SKIP_FILES = {'.DS_Store', 'Thumbs.db', '.gitkeep'}

    def scan_directory(self, directory):
        return list(self._walk_filtered(directory))

    def estimate_folder_size(self, directory: str) -> tuple[int, int, list[str]]:
        total_bytes = 0
        files = []
        for filepath in self._walk_filtered(directory):
            total_bytes += os.path.getsize(filepath)
            files.append(filepath)
        return total_bytes, len(files), files

    def extract_text(self, filepath):
        segments = self.extract_text_with_meta(filepath)
        if not segments:
            return None
        return '\n'.join(text for text, _ in segments)

    def extract_text_with_meta(self, filepath) -> list[tuple[str, dict]]:
        ext = os.path.splitext(filepath)[1].lower()
        if ext == '.pdf':
            return self._extract_pdf_with_meta(filepath)
        if ext in self.OFFICE_EXTENSIONS:
            return self._extract_office_with_meta(ext, filepath)
        if ext in self.TEXT_EXTENSIONS:
            return [(self._extract_text_file(filepath), {'type': 'text'})]
        return []

    def _walk_filtered(self, directory):
        specs = []
        gitignore_path = os.path.join(directory, '.gitignore')
        if os.path.isfile(gitignore_path):
            with open(gitignore_path, 'r') as f:
                specs.append(pathspec.PathSpec.from_lines('gitwildmatch', f))

        for root, dirs, filenames in os.walk(directory):
            dirs[:] = [d for d in dirs if d not in self.SKIP_DIRS]

            nested_gitignore = os.path.join(root, '.gitignore')
            if root != directory and os.path.isfile(nested_gitignore):
                with open(nested_gitignore, 'r') as f:
                    specs.append(pathspec.PathSpec.from_lines('gitwildmatch', f))

            for filename in filenames:
                if filename in self.SKIP_FILES:
                    continue
                filepath = os.path.join(root, filename)
                rel_path = os.path.relpath(filepath, directory)
                if any(spec.match_file(rel_path) for spec in specs):
                    continue
                ext = os.path.splitext(filename)[1].lower()
                if ext in self.SUPPORTED_EXTENSIONS:
                    yield filepath

    def _extract_text_file(self, filepath):
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_pdf_with_meta(self, filepath):
        reader = PdfReader(filepath)
        return [(page.extract_text(), {'page': i + 1}) for i, page in enumerate(reader.pages)]

    def _extract_office_with_meta(self, ext, filepath):
        if ext == '.docx':
            return self._extract_docx_with_meta(filepath)
        if ext == '.xlsx':
            return self._extract_xlsx_with_meta(filepath)
        if ext == '.pptx':
            return self._extract_pptx_with_meta(filepath)
        return []

    def _extract_docx_with_meta(self, filepath):
        from docx import Document
        doc = Document(filepath)
        text = '\n'.join(p.text for p in doc.paragraphs)
        return [(text, {'page': 1})]

    def _extract_xlsx_with_meta(self, filepath):
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True)
        segments = []
        for sheet in wb.worksheets:
            lines = []
            for row in sheet.iter_rows(values_only=True):
                lines.append('\t'.join(str(cell) if cell is not None else '' for cell in row))
            segments.append(('\n'.join(lines), {'sheet': sheet.title}))
        return segments

    def _extract_pptx_with_meta(self, filepath):
        from pptx import Presentation
        prs = Presentation(filepath)
        segments = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            segments.append(('\n'.join(texts), {'slide': i + 1}))
        return segments
